"""Document processing for various file types."""

import os
from typing import List, Dict, Any
from pathlib import Path

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    print("Warning: PyPDF2 not available. PDF processing disabled.")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("Warning: python-docx not available. DOCX processing disabled.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not available. PPTX processing disabled.")
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.schema import Document as LangchainDocument
    HAS_LANGCHAIN = True
except ImportError:
    HAS_LANGCHAIN = False
    print("Warning: LangChain not available. Using simple text splitter.")
    
    # Simple fallback text splitter
    class RecursiveCharacterTextSplitter:
        def __init__(self, chunk_size=1000, chunk_overlap=200, **kwargs):
            self.chunk_size = chunk_size
            self.chunk_overlap = chunk_overlap
        
        def split_text(self, text):
            """Simple text splitting fallback."""
            chunks = []
            for i in range(0, len(text), self.chunk_size - self.chunk_overlap):
                chunk = text[i:i + self.chunk_size]
                if chunk.strip():
                    chunks.append(chunk)
            return chunks
    
    # Simple document class
    class LangchainDocument:
        def __init__(self, page_content, metadata=None):
            self.page_content = page_content
            self.metadata = metadata or {}

from .config import settings


class DocumentProcessor:
    """Process various document types and split them into chunks."""
    
    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        """Initialize the document processor.
        
        Args:
            chunk_size: Size of text chunks (defaults to config setting)
            chunk_overlap: Overlap between chunks (defaults to config setting)
        """
        self.chunk_size = chunk_size or settings.chunk_size
        self.chunk_overlap = chunk_overlap or settings.chunk_overlap
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def process_pdf(self, file_path: str) -> str:
        """Extract text from PDF file.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Extracted text content
        """
        if not HAS_PDF:
            return "PDF processing not available. Please install PyPDF2."
        
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error processing PDF {file_path}: {str(e)}")
            return ""
        
        return text
    
    def process_docx(self, file_path: str) -> str:
        """Extract text from Word document.
        
        Args:
            file_path: Path to the Word document
            
        Returns:
            Extracted text content
        """
        if not HAS_DOCX:
            return "DOCX processing not available. Please install python-docx."
        
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            print(f"Error processing DOCX {file_path}: {str(e)}")
            return ""
        
        return text
    
    def process_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Extracted text content
        """
        if not HAS_PPTX:
            return "PPTX processing not available. Please install python-pptx."
        
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error processing PPTX {file_path}: {str(e)}")
            return ""
        
        return text
    
    def process_txt(self, file_path: str) -> str:
        """Read text from plain text file.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            File content
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error processing TXT {file_path}: {str(e)}")
            return ""
    
    def process_file(self, file_path: str) -> str:
        """Process a file based on its extension.
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            Extracted text content
        """
        file_extension = Path(file_path).suffix.lower()
        
        processors = {
            '.pdf': self.process_pdf,
            '.docx': self.process_docx,
            '.pptx': self.process_pptx,
            '.txt': self.process_txt,
            '.md': self.process_txt
        }
        
        processor = processors.get(file_extension)
        if processor:
            return processor(file_path)
        else:
            print(f"Unsupported file type: {file_extension}")
            return ""
    
    def process_directory(self, directory_path: str) -> List[LangchainDocument]:
        """Process all supported files in a directory.
        
        Args:
            directory_path: Path to the directory containing documents
            
        Returns:
            List of processed document chunks
        """
        documents = []
        directory = Path(directory_path)
        
        if not directory.exists():
            print(f"Directory does not exist: {directory_path}")
            return documents
        
        supported_extensions = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
        
        for file_path in directory.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                print(f"Processing: {file_path}")
                
                text = self.process_file(str(file_path))
                if text.strip():
                    # Split text into chunks
                    chunks = self.text_splitter.split_text(text)
                    
                    # Create document objects with metadata
                    for i, chunk in enumerate(chunks):
                        doc = LangchainDocument(
                            page_content=chunk,
                            metadata={
                                'source': str(file_path),
                                'chunk_index': i,
                                'total_chunks': len(chunks),
                                'file_type': file_path.suffix.lower()
                            }
                        )
                        documents.append(doc)
        
        print(f"Processed {len(documents)} document chunks from {directory_path}")
        return documents